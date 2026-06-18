from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

NAVY  = RGBColor(0x1F,0x38,0x64)
GRAY  = RGBColor(0x59,0x59,0x59)
WHITE = RGBColor(0xFF,0xFF,0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def title_bar(s, text):
    bar = s.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame; tf.margin_left=Inches(0.5); tf.word_wrap=True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = WHITE
    return s

def add_bullets(s, items, left=Inches(0.6), top=Inches(1.45),
                width=Inches(12.1), height=Inches(5.7), base=18):
    tb = s.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for lvl, text, *opt in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        run = p.add_run(); run.text = ("• " if lvl==0 else "– ") + text
        run.font.size = Pt(base - lvl*2)
        run.font.color.rgb = NAVY if lvl==0 else GRAY
        if opt and opt[0]=="b": run.font.bold = True
        p.space_after = Pt(6)
    return tb

def add_image_fit(s, path, left, top, max_w, max_h):
    iw, ih = Image.open(path).size
    ratio = iw/ih
    w = max_w; h = Emu(int(int(w)/ratio))
    if int(h) > int(max_h):
        h = max_h; w = Emu(int(int(h)*ratio))
    cx = left + Emu(int((int(max_w) - int(w))/2))
    cy = top  + Emu(int((int(max_h) - int(h))/2))
    s.shapes.add_picture(path, cx, cy, width=w, height=h)

def caption(s, text, top):
    tb = s.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(0.4))
    p = tb.text_frame.paragraphs[0]; p.text = text
    p.font.size = Pt(12); p.font.italic = True; p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

A = "assets"

# ---------- 1: Title ----------
s = slide()
bg = s.shapes.add_shape(1,0,0,SW,SH); bg.fill.solid(); bg.fill.fore_color.rgb=NAVY; bg.line.fill.background()
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.6))
tf = tb.text_frame; tf.word_wrap=True
p=tf.paragraphs[0]; p.text="Klassifikation deutscher Verkehrsschilder"
p.font.size=Pt(40); p.font.bold=True; p.font.color.rgb=WHITE
p2=tf.add_paragraph(); p2.text="mit neuronalen Netzen (GTSRB-Datensatz)"
p2.font.size=Pt(24); p2.font.color.rgb=RGBColor(0xC9,0xD3,0xE6)
p3=tf.add_paragraph(); p3.text=""
p4=tf.add_paragraph(); p4.text="KI-Praktikum 2  ·  Name(n)  ·  18.06.2026"
p4.font.size=Pt(15); p4.font.color.rgb=RGBColor(0x9F,0xB0,0xD0)

# ---------- 2: Aufgabe & Ziel ----------
s=title_bar(slide(),"Aufgabenstellung & Ziel")
add_bullets(s,[
 (0,"Ziel: aus einem Bild das Verkehrsschild klassifizieren (43 Klassen)","b"),
 (0,"Drei Teilaufgaben:"),
 (1,"Eigene CNN-Modelle entwickeln (mind. 3 Architekturen)"),
 (1,"Vortrainiertes YOLO-Modell anwenden"),
 (1,"YOLO per Finetuning an GTSRB anpassen"),
 (0,"Vorgaben:"),
 (1,"Jedes eigene Modell ≥ 90 % Validierungs-Accuracy"),
 (1,"Bestes Modell anhand der Konfusionsmatrix auswählen & begründen"),
])

# ---------- 3: Datensatz ----------
s=title_bar(slide(),"Der Datensatz (GTSRB)")
add_bullets(s,[
 (0,"~50.000 Bilder, 43 Klassen von Verkehrsschildern","b"),
 (0,"Einheitlich auf 32 × 32 Pixel (RGB) skaliert"),
 (0,"34.799 Trainings- / 4.410 Validierungs- / 12.630 Testbilder"),
 (0,"Unbalanciert: seltene Klassen nur ~30–60 Bilder, häufige > 200","b"),
 (0,"Vorverarbeitung: Pixel auf [0,1] normiert, Daten gemischt (shuffle)"),
])

# ---------- 4: Methodik ----------
s=title_bar(slide(),"Methodik")
add_bullets(s,[
 (0,"Framework: TensorFlow / Keras (Sequential API)","b"),
 (0,"Optimizer: Adam  ·  Loss: sparse categorical crossentropy  ·  Batch: 64"),
 (0,"EarlyStopping (patience = 3, restore_best_weights) gegen Overfitting","b"),
 (0,"Bewertung: Validierungs-Accuracy + Konfusionsmatrix"),
])

# ---------- 5: 3 Architekturen ----------
s=title_bar(slide(),"Eigene Modelle: 3 Architekturen")
add_bullets(s,[
 (0,"Modell 1 – Basis-CNN","b"),
 (1,"2 Faltungsblöcke (32→64), MaxPool + Dropout, Dense(256) — einfache Referenz"),
 (0,"Modell 2 – Tiefes CNN + Batch Normalization (VGG-artig)","b"),
 (1,"3 Blöcke à 2 Conv (32→64→128), BatchNorm — höchste Kapazität"),
 (0,"Modell 3 – Schlankes, LeNet-inspiriertes CNN","b"),
 (1,"2 Conv (5×5, 16→32), AveragePooling, Dense(120→84) — sehr wenige Parameter"),
 (0,"→ bewusst drei verschiedene Designphilosophien"),
])

# ---------- 6: Ergebnisse eigene Modelle ----------
s=title_bar(slide(),"Ergebnisse der eigenen Modelle (Validierung)")
add_bullets(s,[
 (0,"Modell 1 (Basis-CNN): 96,94 %","b"),
 (0,"Modell 2 (BatchNorm): 97,82 %  ⭐ bestes Modell","b"),
 (0,"Modell 3 (LeNet): 93,47 %","b"),
 (0,"Alle Modelle ≥ 90 % erreicht ✓"),
 (1,"Modell 3 trainiert ~10× schneller, aber ~4 Punkte schwächer"),
], left=Inches(0.6), top=Inches(1.5), width=Inches(6.2), height=Inches(5.5), base=18)
add_image_fit(s, f"{A}/acc_models.png", Inches(6.9), Inches(1.5), Inches(6.1), Inches(5.4))

# ---------- 7: Modellvergleich CM ----------
s=title_bar(slide(),"Modellvergleich per Konfusionsmatrix")
add_image_fit(s, f"{A}/cm_compare.png", Inches(0.4), Inches(1.55), Inches(12.5), Inches(4.5))
caption(s,"Konfusionsmatrizen (Validierung) — bestes Modell: Modell 2 (BatchNorm) mit sauberster Diagonale", Inches(6.25))

# ---------- 8: CM-Analyse ----------
s=title_bar(slide(),"Analyse der Konfusionsmatrix (bestes Modell)")
add_bullets(s,[
 (0,"Häufigste Verwechslungen (echt → vorhergesagt):","b"),
 (1,"Verbot LKW>3,5t → Ende Überholverbot: 25×"),
 (1,"Tempo 120 → Tempo 100: 13×  (ähnliche Ziffern)"),
 (1,"Verengung rechts / Fußgänger → Gefahr allgemein"),
 (0,"Schwächste Klassen: seltene Schilder","b"),
 (1,"Verbot LKW>3,5t 58 %, Verengung rechts 70 %, Fußgänger 77 %"),
 (0,"Zuverlässig (100 %): Vorfahrtstraße, Vorfahrt gewähren, Tempo 80/100","b"),
 (0,"Höchste Accuracy = überzeugendste Matrix? Ja, Modell 2","b"),
], left=Inches(0.6), top=Inches(1.45), width=Inches(7.0), height=Inches(5.7), base=16)
add_image_fit(s, f"{A}/cm_best.png", Inches(7.7), Inches(1.5), Inches(5.4), Inches(5.4))

# ---------- 9: Finales Modell auf Testdaten (ECHTE DATEN) ----------
s=title_bar(slide(),"Finales Modell auf Testdaten (Modell 2)")
add_bullets(s,[
 (0,"Test-Accuracy: 97,65 % auf 12.630 Testbildern","b"),
 (0,"Validierung war 97,82 %  →  nur −0,17 Punkte","b"),
 (1,"sehr gute Generalisierung, praktisch kein Overfitting"),
 (0,"Häufigste Verwechslungen (Test):"),
 (1,"Doppelkurve → Vorfahrtstraße, Tempo 120 → 60, Einfahrt verboten → Vorfahrtstraße"),
 (0,"Schwächste Klassen: Fußgänger (50 %), Doppelkurve (64 %)"),
 (1,"wieder seltene Klassen — gleiches Muster wie bei der Validierung"),
], left=Inches(0.6), top=Inches(1.45), width=Inches(7.0), height=Inches(5.7), base=16)
add_image_fit(s, f"{A}/cm_test.png", Inches(7.7), Inches(1.5), Inches(5.4), Inches(5.4))

# ---------- 10: YOLO pretrained ----------
s=title_bar(slide(),"Vortrainiertes YOLO (ohne Anpassung)")
add_bullets(s,[
 (0,"Modell yolo26n-cls.pt (Ultralytics, auf ImageNet vortrainiert)","b"),
 (0,"Alle 4.410 Validierungsbilder klassifiziert"),
 (0,"Top-erkannte Klassen:","b"),
 (1,"loupe (1227), punching_bag (279), digital_clock (243), milk_can (218) …"),
 (0,"Fazit: völlig ungeeignet — kennt keine Verkehrsschilder,","b"),
 (1,"ordnet nur ImageNet-Alltagsobjekte zu"),
])

# ---------- 11: YOLO Finetuning (BEIDE VARIANTEN) ----------
s=title_bar(slide(),"YOLO mit Finetuning: eingefroren vs. voll")
add_bullets(s,[
 (0,"Auflösung von imgsz 32 auf 64 erhöht — größter Hebel","b"),
 (1,"vorher (imgsz 32): nur ~46 % Top-1"),
 (0,"Variante A – Backbone eingefroren (aufgaben-konform):","b"),
 (1,"Top-1 = 70,41 %  —  Einfrieren begrenzt die Genauigkeit"),
 (0,"Variante B – volles Finetuning (Backbone mit-trainiert):","b"),
 (1,"Top-1 = 95,31 %  ·  Top-5 = 99,71 %  ✓ ≥ 95 % erreicht"),
 (0,"Training auf Apple-GPU (MPS), je 30 Epochen, EarlyStopping"),
], left=Inches(0.6), top=Inches(1.45), width=Inches(6.6), height=Inches(5.7), base=16)
add_image_fit(s, f"{A}/yolo_compare.png", Inches(7.3), Inches(1.6), Inches(5.8), Inches(5.2))

# ---------- 12: YOLO Konfusionsmatrix (NEU) ----------
s=title_bar(slide(),"YOLO (voll finetuned) – Konfusionsmatrix")
add_image_fit(s, f"{A}/yolo_cm.png", Inches(0.4), Inches(1.5), Inches(12.5), Inches(5.2))
caption(s,"Normalisierte Konfusionsmatrix des voll finetunten YOLO-Modells (Validierung, 95,31 % Top-1)", Inches(6.85))

# ---------- 13: Gesamtvergleich ----------
s=title_bar(slide(),"Gesamtvergleich der Ansätze")
add_bullets(s,[
 (0,"Eigenes CNN (Modell 2): 97,82 % Val / 97,65 % Test","b"),
 (0,"YOLO finetuned (voll): 95,31 %","b"),
 (0,"YOLO finetuned (eingefroren): 70,41 %","b"),
 (0,"YOLO vortrainiert: unbrauchbar","b"),
 (0,"Eigenes CNN bleibt vorn — volles YOLO-Finetuning ist aber konkurrenzfähig"),
], left=Inches(0.6), top=Inches(1.5), width=Inches(6.2), height=Inches(5.5), base=18)
add_image_fit(s, f"{A}/acc_overall.png", Inches(6.9), Inches(1.5), Inches(6.1), Inches(5.4))

# ---------- 14: Fazit ----------
s=title_bar(slide(),"Fazit & Ausblick")
add_bullets(s,[
 (0,"3 eigene CNNs trainiert — alle > 90 %, bestes 97,82 % (BatchNorm)","b"),
 (0,"Finales Modell auf Testdaten bestätigt: 97,65 %, kein Overfitting","b"),
 (0,"Vortrainiertes YOLO ohne Anpassung wertlos für Spezialdomäne","b"),
 (0,"YOLO-Finetuning: Auflösung & volles Training entscheidend","b"),
 (1,"eingefroren ~70 %, voll finetuned 95,31 %"),
 (0,"Schlüsselkonzepte: Transfer Learning, Layer einfrieren, BatchNorm,"),
 (1,"Datenungleichgewicht, Konfusionsmatrix-Analyse"),
])

out="KI_Praktikum2_Praesentation.pptx"
prs.save(out)
print("SAVED:", out, "| Folien:", len(prs.slides._sldIdLst))
